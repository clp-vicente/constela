import sys
import json
import io
import base64
import mimetypes
import fitz
import docx
from pptx import Presentation

from gemini_client import GeminiClient

API_KEY = ""

# extraccion de texto

def extract_text_from_pdf(file_path):
    try:
        with fitz.open(file_path) as doc:
            text = "".join(page.get_text() for page in doc)
        return text
    except Exception as e:
        return f"Error al leer PDF: {e}"

def extract_text_from_docx(file_path):
    try:
        doc = docx.Document(file_path)
        text = "\n".join(para.text for para in doc.paragraphs)
        return text
    except Exception as e:
        return f"Error al leer DOCX: {e}"

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error al leer PPTX: {e}"

def send_to_stdout(data_dict):
    try:
        message = json.dumps(data_dict)
        sys.stdout.buffer.write(message.encode('utf-8'))
        sys.stdout.buffer.write(b'\n')
        sys.stdout.flush()
    except Exception as e:
        fallback_error = json.dumps({"error": f"Error interno en el backend al enviar datos: {str(e)}"})
        sys.stdout.buffer.write(fallback_error.encode('utf-8'))
        sys.stdout.buffer.write(b'\n')
        sys.stdout.flush()

def main():
    sys.stdin = io.TextIOWrapper(sys.stdin.buffer, encoding='utf-8')

    model_name = sys.argv[1] if len(sys.argv) > 1 else 'gemini-2.5-flash'
    
    gemini_client = GeminiClient(api_key=API_KEY, model_name=model_name)
    
    if not gemini_client.model:
        send_to_stdout({"error": "API Key inválida o modelo incorrecto."})
        sys.exit(1)

    send_to_stdout({"status": "ready"})

    while True:
        try:
            line = sys.stdin.readline()
            if not line:
                break
            
            request = json.loads(line)

            action = request.get("action")
            if action == 'reset':
                gemini_client.start_new_chat()
                continue
            
            elif action == 'load_history':
                history_from_js = request.get("history", [])
                
                python_compatible_history = []
                for message in history_from_js:
                    if "role" in message and "parts" in message and isinstance(message["parts"], list):
                        transformed_parts = [part.get("text") for part in message["parts"] if "text" in part]
                        python_compatible_history.append({
                            "role": message["role"],
                            "parts": transformed_parts
                        })
                
                gemini_client.load_chat_history(python_compatible_history)
                continue
            
            user_input = request.get("prompt")
            file_path = request.get("filePath")
            
            content_parts = []
            
            if file_path:
                try:
                    file_content_text = ""
                    lower_file_path = file_path.lower()
                    mime_type, _ = mimetypes.guess_type(file_path)

                    if mime_type and mime_type.startswith('image/'):
                        with open(file_path, "rb") as f:
                            encoded_data = base64.b64encode(f.read()).decode('utf-8')
                        
                        image_part = {
                            "inline_data": {
                                "mime_type": mime_type,
                                "data": encoded_data
                            }
                        }
                        content_parts.insert(0, image_part)
                    
                    elif lower_file_path.endswith('.pdf'):
                        file_content_text = extract_text_from_pdf(file_path)
                    
                    elif lower_file_path.endswith('.docx'):
                        file_content_text = extract_text_from_docx(file_path)
                    
                    elif lower_file_path.endswith('.pptx'):
                        file_content_text = extract_text_from_pptx(file_path)

                    if file_content_text:
                        formatted_content = f"--- INICIO DEL CONTENIDO DEL ARCHIVO ADJUNTO ---\n{file_content_text}\n--- FIN DEL CONTENIDO DEL ARCHIVO ADJUNTO ---"
                        content_parts.append(formatted_content)

                except FileNotFoundError:
                    send_to_stdout({"error": f"Archivo no encontrado: {file_path}"})
                    continue
                except Exception as e:
                    send_to_stdout({"error": f"Error al procesar el archivo: {str(e)}"})
                    continue
            
            if user_input:
                content_parts.append(user_input)
            
            if content_parts:
                response = gemini_client.send_message(content_parts)
                send_to_stdout(response)

        except Exception as e:
            send_to_stdout({"error": str(e)})

if __name__ == "__main__":
    main()